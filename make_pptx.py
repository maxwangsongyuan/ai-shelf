#!/usr/bin/env python3
"""Generate The AI Shelf hackathon presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
DARK_BG = RGBColor(0x12, 0x12, 0x12)
CARD_BG = RGBColor(0x1E, 0x1E, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0x9C, 0x9C)
LIGHT_GRAY = RGBColor(0xD4, 0xD4, 0xD4)
YELLOW = RGBColor(0xFA, 0xCC, 0x15)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
CYAN = RGBColor(0x06, 0xB6, 0xD4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color=BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    adj = shape.adjustments
    if len(adj) > 0:
        adj[0] = 0.05
    return shape

def add_multiline_box(slide, left, top, width, height, lines, default_size=16, default_color=WHITE):
    """lines: list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Calibri'
        p.alignment = align
        p.space_after = Pt(4)
    return txBox

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_text_box(slide, 0, 1.2, 13.333, 1.2, "The AI Shelf", 60, YELLOW, True, PP_ALIGN.CENTER)
add_text_box(slide, 0, 2.5, 13.333, 0.8, "Measure Your Brand's Visibility on the Invisible AI Shelf", 24, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Shelf metaphor explanation
card = add_rounded_rect(slide, 2.5, 3.5, 8.333, 1.2, CARD_BG)
add_multiline_box(slide, 2.8, 3.55, 7.8, 1.1, [
    ("Think of AI as a store shelf. When someone asks ChatGPT", 16, WHITE, False, PP_ALIGN.CENTER),
    ('"best dental lab in NJ", the brands it mentions are on that shelf.', 16, WHITE, False, PP_ALIGN.CENTER),
    ("If you're not mentioned — you don't exist.", 16, YELLOW, True, PP_ALIGN.CENTER),
])

add_text_box(slide, 0, 5.0, 13.333, 0.5, "Pulse AI x CASTGNY Hackathon  |  Max Wang", 18, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 0, 5.6, 13.333, 0.5, "ai-shelf.vercel.app", 16, BLUE, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.8, 0.4, 12, 0.8, "The Invisible Crisis", 40, WHITE, True)
add_text_box(slide, 0.8, 1.1, 10, 0.5, "What happens after a sales rep leaves a brochure?", 18, GRAY)

# Story flow - left side
card = add_rounded_rect(slide, 0.8, 1.8, 5.5, 5.0)
add_multiline_box(slide, 1.1, 1.9, 5.0, 4.8, [
    ("The Story of iCo Dental Group", 20, YELLOW, True),
    ("", 8),
    ("Real NJ dental lab. 500+ practices. 5 years in business.", 16, LIGHT_GRAY),
    ("", 8),
    ("Their sales reps visit dentist offices every day,", 16, WHITE),
    ("leaving brochures and price quotes.", 16, WHITE),
    ("", 8),
    ("But after the sales rep leaves...", 16, GRAY),
    ("", 8),
    ("The dentist picks up their phone and asks AI:", 16, WHITE),
    ("", 6),
    ('"Is iCo Dental Group any good?"', 20, CYAN, True),
    ("", 8),
    ("AI responds:", 14, GRAY),
    ('"I don\'t have specific information', 18, RED, True),
    ('  about iCo Dental Group..."', 18, RED, True),
])

# Right side - consequence
card2 = add_rounded_rect(slide, 6.8, 1.8, 5.7, 2.2)
add_multiline_box(slide, 7.1, 1.9, 5.2, 2.0, [
    ("The Result", 20, RED, True),
    ("", 6),
    ("Trust destroyed in 10 seconds.", 18, WHITE),
    ("Brochure goes in the trash.", 18, WHITE),
    ("5 years of marketing — invisible to AI.", 18, GRAY),
])

# GEO context
card3 = add_rounded_rect(slide, 6.8, 4.3, 5.7, 2.5)
add_multiline_box(slide, 7.1, 4.4, 5.2, 2.3, [
    ("GEO: The Next SEO", 20, YELLOW, True),
    ("", 6),
    ("SEO = rank on Google (last 10 years)", 15, GRAY),
    ("GEO = appear in AI answers (now)", 15, WHITE, True),
    ("", 6),
    ("24+ companies, $200M+ raised globally", 15, GRAY),
    ("Princeton research: GEO can boost", 15, GRAY),
    ("AI brand visibility by 40%", 15, GREEN, True),
])

# ============================================================
# SLIDE 3: What We Built
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0, 0.4, 13.333, 0.8, "What The AI Shelf Does", 40, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 0, 1.1, 13.333, 0.5, "Three steps: Measure, Diagnose, Prescribe", 20, GRAY, False, PP_ALIGN.CENTER)

# Three cards
labels = [
    ("Measure", BLUE, "Where are you on the\nAI shelf?", "Share of Model Score:\nDiscovery Score + Trust Score\ntell you exactly how visible\nyour brand is in AI responses."),
    ("Diagnose", YELLOW, "Why are you there?", "Gap Analysis compares your\nbrand vs competitors across\nWikipedia, forums, reviews,\nstructured data, and more."),
    ("Prescribe", GREEN, "How to move up?", "Prioritized Action Plan with\nauto-generated content:\nWikipedia drafts, Schema.org\ncode, comparison pages."),
]

for i, (title, color, subtitle, desc) in enumerate(labels):
    x = 1.0 + i * 4.0
    card = add_rounded_rect(slide, x, 1.9, 3.5, 4.5)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.3), Inches(2.1), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(28)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + 0.2, 3.2, 3.1, 0.5, title, 26, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 3.7, 3.1, 0.6, subtitle, 16, WHITE, False, PP_ALIGN.CENTER)
    add_multiline_box(slide, x + 0.3, 4.4, 2.9, 2.0, [
        (desc, 14, LIGHT_GRAY),
    ])

# ============================================================
# SLIDE 4: Two Scenarios
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0, 0.4, 13.333, 0.8, "Dual-Scenario Analysis", 40, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 0, 1.1, 13.333, 0.5, "We cover BOTH ways customers interact with AI about your brand", 18, GRAY, False, PP_ALIGN.CENTER)

# Discovery card
card = add_rounded_rect(slide, 0.8, 1.9, 5.5, 4.8)
add_multiline_box(slide, 1.1, 2.0, 5.0, 4.6, [
    ("Discovery Search", 24, BLUE, True),
    ("Can AI find you without knowing your name?", 14, GRAY),
    ("", 8),
    ("Example queries:", 16, WHITE, True),
    ('"best dental lab in New Jersey"', 15, CYAN),
    ('"affordable zirconia crown lab near me"', 15, CYAN),
    ('"fast turnaround dental lab NJ"', 15, CYAN),
    ("", 8),
    ("iCo Result: 0% Discovery Score", 18, RED, True),
    ("Not mentioned in ANY discovery query.", 14, GRAY),
    ("", 6),
    ("= Invisible on the AI shelf", 16, RED),
])

# Verification card
card = add_rounded_rect(slide, 7.0, 1.9, 5.5, 4.8)
add_multiline_box(slide, 7.3, 2.0, 5.0, 4.6, [
    ("Verification Search", 24, PURPLE, True),
    ("Can AI vouch for you when asked directly?", 14, GRAY),
    ("", 8),
    ("Example queries:", 16, WHITE, True),
    ('"Is iCo Dental Group reputable?"', 15, CYAN),
    ('"iCo Dental Group services and prices?"', 15, CYAN),
    ('"Reviews for iCo Dental Group?"', 15, CYAN),
    ("", 8),
    ("iCo Result: 67% Trust Score", 18, YELLOW, True),
    ("AI knows the name but can't describe it.", 14, GRAY),
    ("", 6),
    ("= Known but not trusted", 16, YELLOW),
])

# Key insight at bottom
add_text_box(slide, 0, 6.8, 13.333, 0.5,
    'Key Insight: "AI knows you, but won\'t recommend you."', 18, YELLOW, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: How It Works - 6-Step Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0, 0.3, 13.333, 0.7, "6-Step Analysis Pipeline", 36, WHITE, True, PP_ALIGN.CENTER)

steps = [
    ("1", "Query\nGeneration", "Claude generates 6\nhigh-intent queries\n(3 discovery + 3 verification)\nwith intent analysis", BLUE),
    ("2", "AI\nResponses", "Claude answers each\nquery. JS string matching\ndetects brand mentions,\nsentiment, competitors", BLUE),
    ("3", "Web Reality\nCheck", "Tavily searches same\n6 queries + brand\npresence. Compares\nAI vs real web data", CYAN),
    ("4", "Scoring", "Local computation.\nDiscovery Score +\nTrust Score + Radar.\nArchetype classification", GREEN),
    ("5", "Gap\nAnalysis", "Tavily searches brand\n+ competitor presence.\nClaude finds specific\ngaps with evidence", YELLOW),
    ("6", "Action\nPlan", "Claude generates\nprioritized actions with\nauto-generated content\n(Wiki, Schema, etc.)", PURPLE),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.5 + i * 2.1
    # Card
    card = add_rounded_rect(slide, x, 1.2, 1.9, 5.5)
    # Step number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.55), Inches(1.4), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + 0.1, 2.4, 1.7, 0.9, title, 16, color, True, PP_ALIGN.CENTER)
    add_multiline_box(slide, x + 0.15, 3.4, 1.6, 3.0, [
        (desc, 12, LIGHT_GRAY),
    ])

    # Arrow between steps
    if i < 5:
        arrow_x = x + 1.9
        add_text_box(slide, arrow_x, 3.5, 0.3, 0.5, ">", 20, GRAY, True, PP_ALIGN.CENTER)

# API call summary at bottom
add_rounded_rect(slide, 0.5, 6.8, 12.3, 0.6, CARD_BG)
add_text_box(slide, 0.8, 6.85, 12, 0.5,
    "Per analysis:  6x Claude calls  +  12x Tavily calls  |  SSE streaming  |  ~30s total  |  ~$0.05 cost",
    14, GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: Technical Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0, 0.3, 13.333, 0.7, "Technical Architecture", 36, WHITE, True, PP_ALIGN.CENTER)

# Tech stack card
card = add_rounded_rect(slide, 0.8, 1.2, 5.5, 5.8)
add_multiline_box(slide, 1.1, 1.3, 5.0, 5.6, [
    ("Tech Stack", 22, YELLOW, True),
    ("", 8),
    ("Frontend", 16, BLUE, True),
    ("Next.js 15 App Router + TypeScript", 14, WHITE),
    ("Tailwind CSS + shadcn/ui", 14, WHITE),
    ("Recharts (radar chart)", 14, WHITE),
    ("", 8),
    ("Backend", 16, GREEN, True),
    ("Next.js API Route (Node.js runtime)", 14, WHITE),
    ("SSE (Server-Sent Events) streaming", 14, WHITE),
    ("Zod schema + generateObject", 14, WHITE),
    ("", 8),
    ("AI & Search", 16, PURPLE, True),
    ("Claude Haiku 4.5 (via Vercel AI SDK)", 14, WHITE),
    ("Tavily Search API", 14, WHITE),
    ("JS string matching (brand analysis)", 14, WHITE),
    ("", 8),
    ("Deployment", 16, CYAN, True),
    ("Vercel (Hobby plan, free)", 14, WHITE),
    ("60s max serverless function duration", 14, WHITE),
])

# Data flow card
card = add_rounded_rect(slide, 6.8, 1.2, 5.7, 5.8)
add_multiline_box(slide, 7.1, 1.3, 5.2, 5.6, [
    ("Data Flow", 22, YELLOW, True),
    ("", 8),
    ("User Input", 14, GRAY),
    ("  Brand + Category + Region", 15, WHITE),
    ("        |", 14, GRAY),
    ("  POST /api/analyze (SSE stream)", 14, CYAN),
    ("        |", 14, GRAY),
    ("Step 1: Claude generateObject", 14, BLUE),
    ("  -> 6 queries with intent", 13, LIGHT_GRAY),
    ("        |", 14, GRAY),
    ("Step 2+3: Parallel execution", 14, BLUE),
    ("  Claude x6 + Tavily x6", 13, LIGHT_GRAY),
    ("        |", 14, GRAY),
    ("Step 4: Local score computation", 14, GREEN),
    ("  + optional radar (Claude)", 13, LIGHT_GRAY),
    ("        |", 14, GRAY),
    ("Step 5: Tavily x6 + Claude analysis", 14, YELLOW),
    ("  Brand vs competitor presence", 13, LIGHT_GRAY),
    ("        |", 14, GRAY),
    ("Step 6: Claude generateObject", 14, PURPLE),
    ("  -> Prioritized action plan", 13, LIGHT_GRAY),
    ("        |", 14, GRAY),
    ("Frontend renders each step live", 14, WHITE, True),
])

# ============================================================
# SLIDE 7: Live Results - Score
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0, 0.3, 13.333, 0.7, "Live Results: iCo Dental Group", 36, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 0, 0.9, 13.333, 0.5, "dental lab  |  New Jersey  |  Real brand, real data", 16, GRAY, False, PP_ALIGN.CENTER)

# Big score
card = add_rounded_rect(slide, 0.8, 1.6, 3.5, 4.0)
add_text_box(slide, 0.8, 1.8, 3.5, 1.5, "34", 80, YELLOW, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 3.4, 3.5, 0.5, "/100", 24, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 4.0, 3.5, 0.5, "Niche Player", 20, YELLOW, True, PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 4.6, 3.5, 0.5, "AI knows you, won't recommend", 12, GRAY, False, PP_ALIGN.CENTER)

# Discovery vs Trust
card = add_rounded_rect(slide, 4.8, 1.6, 3.8, 1.8)
add_multiline_box(slide, 5.1, 1.7, 3.3, 1.6, [
    ("Discovery Score", 18, BLUE, True),
    ("0%", 36, RED, True),
    ("Not in ANY general search", 13, GRAY),
])

card = add_rounded_rect(slide, 4.8, 3.8, 3.8, 1.8)
add_multiline_box(slide, 5.1, 3.9, 3.3, 1.6, [
    ("Trust Score", 18, PURPLE, True),
    ("67%", 36, YELLOW, True),
    ("Mentioned when asked directly", 13, GRAY),
])

# Query results breakdown
card = add_rounded_rect(slide, 9.1, 1.6, 3.5, 4.0)
add_multiline_box(slide, 9.4, 1.7, 3.0, 3.8, [
    ("Query Results  (3/6 mentioned)", 14, WHITE, True),
    ("", 6),
    ("Discovery Queries", 13, BLUE, True),
    ("  best dental lab NJ", 12, GRAY),
    ("  -> Not Mentioned", 12, RED),
    ("  reliable lab fast turnaround", 12, GRAY),
    ("  -> Not Mentioned", 12, RED),
    ("  dentures and prosthetics", 12, GRAY),
    ("  -> Not Mentioned", 12, RED),
    ("", 4),
    ("Verification Queries", 13, PURPLE, True),
    ("  Is iCo reputable?", 12, GRAY),
    ("  -> Mentioned (Neutral)", 12, GREEN),
    ("  iCo services & prices?", 12, GRAY),
    ("  -> Mentioned (Negative)", 12, YELLOW),
    ("  iCo reviews & feedback?", 12, GRAY),
    ("  -> Mentioned (Positive)", 12, GREEN),
])

# Radar chart note
card = add_rounded_rect(slide, 0.8, 5.8, 11.8, 1.0)
add_multiline_box(slide, 1.1, 5.85, 11.3, 0.9, [
    ("5-Dimension Radar:  Recognition | Sentiment | Relevance | Citation | Competitive", 14, GRAY),
    ("Interactive radar chart available in live demo at ai-shelf.vercel.app", 13, BLUE),
])

# ============================================================
# SLIDE 8: Gap Analysis + Action Plan
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0, 0.3, 13.333, 0.7, "Gap Analysis & Action Plan", 36, WHITE, True, PP_ALIGN.CENTER)

# Gap analysis card
card = add_rounded_rect(slide, 0.8, 1.2, 5.5, 5.8)
add_multiline_box(slide, 1.1, 1.3, 5.0, 5.6, [
    ("Why AI Doesn't Know iCo", 20, RED, True),
    ("(Gap Analysis — with evidence links)", 13, GRAY),
    ("", 8),
    ("Wikipedia Page", 15, WHITE, True),
    ("  iCo: no     Competitor: no", 13, GRAY),
    ("", 4),
    ("Industry Association (ICOI)", 15, WHITE, True),
    ("  iCo: no     Competitor: no", 13, GRAY),
    ("", 4),
    ("Google Business Profile", 15, WHITE, True),
    ("  iCo: YES    Competitor: no", 13, GREEN),
    ("", 4),
    ("LinkedIn Company Page", 15, WHITE, True),
    ("  iCo: YES    Competitor: no", 13, GREEN),
    ("", 4),
    ("Domain Authority", 15, WHITE, True),
    ("  Fragmented across multiple domains", 13, RED),
    ("", 4),
    ("Content Depth", 15, WHITE, True),
    ("  Only social snippets, no rich content", 13, RED),
    ("", 4),
    ("Brand Name Disambiguation", 15, WHITE, True),
    ("  'iCo' matches unrelated entities", 13, RED),
    ("", 6),
    ("Every gap has a clickable evidence URL", 13, CYAN),
])

# Action plan card
card = add_rounded_rect(slide, 6.8, 1.2, 5.7, 5.8)
add_multiline_box(slide, 7.1, 1.3, 5.2, 5.6, [
    ("Prioritized Action Plan", 20, GREEN, True),
    ("(Auto-generated by Claude)", 13, GRAY),
    ("", 8),
    ("Priority 1: This Week", 16, RED, True),
    ("  Schema.org JSON-LD markup", 13, WHITE),
    ("    -> +40-60% entity recognition", 12, GREEN),
    ("  Service taxonomy & content hub", 13, WHITE),
    ("    -> +200-300% discoverable content", 12, GREEN),
    ("  Domain consolidation", 13, WHITE),
    ("    -> -70% entity confusion", 12, GREEN),
    ("", 6),
    ("Priority 2: This Month", 16, YELLOW, True),
    ("  ICOI membership application", 13, WHITE),
    ("  Google Business Profile optimization", 13, WHITE),
    ("  LinkedIn thought leadership strategy", 13, WHITE),
    ("", 6),
    ("Priority 3: Ongoing", 16, BLUE, True),
    ("  Wikipedia disambiguation strategy", 13, WHITE),
    ("  Industry publication mentions", 13, WHITE),
    ("  Backlink strategy", 13, WHITE),
    ("", 8),
    ("Each action includes auto-generated", 13, CYAN),
    ("content samples ready to deploy", 13, CYAN),
])

# ============================================================
# SLIDE 9: Hackathon Requirements
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0, 0.3, 13.333, 0.7, "Hackathon Requirements", 36, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 0, 0.9, 13.333, 0.5, "How The AI Shelf meets every requirement", 18, GRAY, False, PP_ALIGN.CENTER)

requirements = [
    ("5+ high-intent queries", "6 queries (3 discovery + 3 verification)\nwith intent analysis for each", GREEN),
    ("2+ AI systems / search tools", "Claude (AI responses) + Tavily (web search)\n= 2 distinct systems", GREEN),
    ("Extract brand mentions,\ncompetitors, position, sentiment", "JS string matching extracts all four\nfrom each AI response automatically", GREEN),
    ("Calculate Share of Model", "Discovery Score + Trust Score\n+ Overall Score + Archetype badge", GREEN),
    ("Show evidence", "Every step expandable with raw data\n+ clickable source URLs", GREEN),
    ("Competitor analysis", "Gap Analysis table: brand vs competitor\nacross 7+ factors with evidence", GREEN),
    ("Auto-generate improvement content", "Wikipedia drafts, Schema.org code,\ncomparison page outlines — ready to use", GREEN),
    ("Prioritized action plan", "3-tier priority: This Week / Month / Ongoing\nwith expected impact metrics", GREEN),
]

for i, (req, how, color) in enumerate(requirements):
    row = i // 2
    col = i % 2
    x = 0.8 + col * 6.2
    y = 1.5 + row * 1.4
    card = add_rounded_rect(slide, x, y, 5.8, 1.2)
    # Checkmark
    add_text_box(slide, x + 0.15, y + 0.15, 0.4, 0.4, "OK", 12, GREEN, True)
    add_multiline_box(slide, x + 0.5, y + 0.1, 2.2, 1.0, [
        (req, 13, WHITE, True),
    ])
    add_multiline_box(slide, x + 2.8, y + 0.1, 2.8, 1.0, [
        (how, 11, LIGHT_GRAY),
    ])

# ============================================================
# SLIDE 10: Why We Win
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0, 0.3, 13.333, 0.7, "Why We Win", 40, YELLOW, True, PP_ALIGN.CENTER)

# Judging criteria
card = add_rounded_rect(slide, 0.8, 1.2, 5.5, 3.5)
add_multiline_box(slide, 1.1, 1.3, 5.0, 3.3, [
    ("Judging Criteria", 20, WHITE, True),
    ("", 6),
    ("Reliability (30%)", 16, GREEN, True),
    ("  Vercel deployed + every conclusion has", 13, LIGHT_GRAY),
    ("  clickable source links + error handling", 13, LIGHT_GRAY),
    ("", 4),
    ("Insight Quality (25%)", 16, BLUE, True),
    ("  Dual-scenario analysis + intent labels", 13, LIGHT_GRAY),
    ("  + Trust Score is our unique metric", 13, LIGHT_GRAY),
    ("", 4),
    ("Technical Robustness (25%)", 16, PURPLE, True),
    ("  Tavily deep usage + structured output", 13, LIGHT_GRAY),
    ("  + parallel API + SSE streaming", 13, LIGHT_GRAY),
    ("", 4),
    ("Clarity to CEO (20%)", 16, YELLOW, True),
    ("  One sentence: \"AI doesn't know you\"", 13, LIGHT_GRAY),
    ("  + 0% vs 67% comparison = instant clarity", 13, LIGHT_GRAY),
])

# Differentiators
card = add_rounded_rect(slide, 6.8, 1.2, 5.7, 3.5)
add_multiline_box(slide, 7.1, 1.3, 5.2, 3.3, [
    ("What Makes Us Different", 20, WHITE, True),
    ("", 6),
    ("Others may do:", 14, GRAY),
    ("  Only discovery analysis", 13, GRAY),
    ("We do:", 14, GREEN, True),
    ("  Discovery + Verification (full journey)", 13, GREEN),
    ("", 4),
    ("Others may do:", 14, GRAY),
    ("  Just keywords", 13, GRAY),
    ("We do:", 14, GREEN, True),
    ("  Each query has intent analysis", 13, GREEN),
    ("", 4),
    ("Others may do:", 14, GRAY),
    ("  One SoM score", 13, GRAY),
    ("We do:", 14, GREEN, True),
    ("  Discovery Score + Trust Score (2D)", 13, GREEN),
    ("", 4),
    ("Others may do:", 14, GRAY),
    ("  Nike / HubSpot demo", 13, GRAY),
    ("We do:", 14, GREEN, True),
    ("  Real dental lab with real pain points", 13, GREEN),
])

# Bottom tagline
card = add_rounded_rect(slide, 0.8, 5.0, 11.7, 2.0)
add_multiline_box(slide, 1.2, 5.1, 11.0, 1.8, [
    ("The Pitch in One Sentence", 18, YELLOW, True),
    ("", 6),
    ('"iCo spent 5 years doing outbound sales, but every visit\'s conversion rate', 16, WHITE),
    ('was silently undermined by AI saying "I don\'t know this company."', 16, WHITE),
    ('In 30 seconds, we found the blind spot, diagnosed the cause, and wrote the prescription.', 16, WHITE),
    ('Any brand, any industry, can run the same check."', 16, WHITE),
])

# ============================================================
# SLIDE 11: Thank You / CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0, 1.5, 13.333, 1.2, "Try It Yourself", 50, YELLOW, True, PP_ALIGN.CENTER)
add_text_box(slide, 0, 2.8, 13.333, 0.8, "ai-shelf.vercel.app", 36, BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, 0, 4.0, 13.333, 0.5, 'Click "Try: iCo Dental Group" and hit Analyze', 20, GRAY, False, PP_ALIGN.CENTER)

add_multiline_box(slide, 0, 5.0, 13.333, 1.5, [
    ("Built with:", 16, GRAY, False, PP_ALIGN.CENTER),
    ("Claude Haiku 4.5  |  Tavily Search  |  Vercel AI SDK  |  Next.js 15", 18, WHITE, True, PP_ALIGN.CENTER),
    ("", 12),
    ("Max Wang  |  Pulse AI x CASTGNY Hackathon  |  2026", 16, GRAY, False, PP_ALIGN.CENTER),
])

# Save
output_path = '/Users/maxwsy/workspace/ai-shelf/The_AI_Shelf_Presentation.pptx'
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
